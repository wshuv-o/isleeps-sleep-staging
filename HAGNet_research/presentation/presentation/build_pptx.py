"""Final iSLEEPS proposal deck, structured to the teacher's Presentation_Guideline.md.

Slide headers 2-7 are the guideline's EXACT section titles. 1 cover + 9 content = 10 slides.
No em-dashes. Output name via env PPTX_OUT (default iSLEEPS_proposal.pptx).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(HERE, "figures")
OUT = os.path.join(HERE, os.environ.get("PPTX_OUT", "iSLEEPS_proposal.pptx"))

NAVY = RGBColor(0x1B, 0x2A, 0x4A); BLUE = RGBColor(0x2C, 0x5F, 0x8A)
TEAL = RGBColor(0x2A, 0x9D, 0x8F); ORANGE = RGBColor(0xE7, 0x6F, 0x51)
GRAY = RGBColor(0x5A, 0x63, 0x72); LGRAY = RGBColor(0xEC, 0xF0, 0xF4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF); DARK = RGBColor(0x22, 0x2A, 0x35)
GREEN = RGBColor(0x2E, 0x7D, 0x57); RED = RGBColor(0xB0, 0x3A, 0x2E)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _set(par, text, size, color, bold=False, italic=False):
    par.text = text; r = par.runs[0]
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = "Calibri"


def box(slide, l, t, w, h, fill=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h); sp.shadow.inherit = False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.fill.background()
    return sp


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    return tf


def footer(slide, n):
    box(slide, 0, SH - Inches(0.32), SW, Inches(0.32), fill=NAVY)
    tf = textbox(slide, Inches(0.3), SH - Inches(0.33), Inches(11.6), Inches(0.3), MSO_ANCHOR.MIDDLE)
    _set(tf.paragraphs[0], "Research Proposal · Lesion-Resolved Sleep Microstructure in Ischemic Stroke (iSLEEPS)", 9, WHITE)
    tn = textbox(slide, SW - Inches(1.0), SH - Inches(0.33), Inches(0.7), Inches(0.3), MSO_ANCHOR.MIDDLE)
    p = tn.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT; _set(p, str(n), 9, WHITE)


def content_slide(n, title, tag=None):
    s = prs.slides.add_slide(BLANK)
    box(s, 0, 0, SW, Inches(1.1), fill=NAVY)
    box(s, 0, Inches(1.1), SW, Inches(0.06), fill=ORANGE)
    tf = textbox(s, Inches(0.5), Inches(0.12), SW - Inches(3.4), Inches(0.96), MSO_ANCHOR.MIDDLE)
    _set(tf.paragraphs[0], title, 23, WHITE, bold=True)
    if tag:
        tg = textbox(s, SW - Inches(3.7), Inches(0.28), Inches(3.4), Inches(0.55), MSO_ANCHOR.MIDDLE)
        p = tg.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT; _set(p, tag, 11.5, RGBColor(0xBF, 0xD3, 0xE6), italic=True)
    footer(s, n)
    return s


def bullets(slide, items, l, t, w, h, size=16, gap=6):
    tf = textbox(slide, l, t, w, h)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        lvl = 0
        if isinstance(it, tuple): it, lvl = it
        p.space_after = Pt(gap); p.level = lvl
        _set(p, ("> " if lvl == 0 else "- ") + it, size if lvl == 0 else size - 2, DARK if lvl == 0 else GRAY)
    return tf


def chip(slide, l, t, w, text, fill, size=13, h=0.55):
    box(slide, l, t, w, Inches(h), fill=fill)
    tf = textbox(slide, l, t, w, Inches(h), MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; _set(p, text, size, WHITE, bold=True)


def pic(slide, path, l, t, w=None, h=None):
    if os.path.exists(path): slide.shapes.add_picture(path, l, t, width=w, height=h)


def refnote(slide, text, y=6.82):
    tf = textbox(slide, Inches(0.6), Inches(y), Inches(12.15), Inches(0.32))
    _set(tf.paragraphs[0], "Ref: " + text, 9.5, GRAY, italic=True)


# ===== 1. COVER =====
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, SW, SH, fill=NAVY)
box(s, Inches(0.9), Inches(0.75), Inches(2.6), Inches(0.5), fill=ORANGE)
tf = textbox(s, Inches(0.9), Inches(0.75), Inches(2.6), Inches(0.5), MSO_ANCHOR.MIDDLE)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; _set(p, "RESEARCH PROPOSAL", 14, WHITE, bold=True)
tf = textbox(s, Inches(0.9), Inches(1.45), SW - Inches(1.8), Inches(1.5), MSO_ANCHOR.BOTTOM)
_set(tf.paragraphs[0], "Lesion-Resolved Sleep Microstructure in Ischemic Stroke", 33, WHITE, bold=True)
tf2 = textbox(s, Inches(0.9), Inches(3.0), SW - Inches(1.8), Inches(1.1))
_set(tf2.paragraphs[0], "Robust, interpretable EEG sleep-stage classification, used as an instrument to map how focal injury reorganizes sleep", 17, RGBColor(0xBF, 0xD3, 0xE6))
box(s, Inches(0.9), Inches(4.35), SW - Inches(1.8), Inches(0.04), fill=ORANGE)
tf3 = textbox(s, Inches(0.9), Inches(4.5), SW - Inches(1.8), Inches(0.95))
for i, line in enumerate([
    "Dataset: iSLEEPS (2026), Polysomnography in 100 Indian Ischemic Stroke Patients",
    "Supervised 5-class sleep-stage classification; staging as the enabling instrument"]):
    p = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
    _set(p, line, 13.5, WHITE); p.space_after = Pt(3)
# presenters
pl = textbox(s, Inches(0.9), Inches(5.55), SW - Inches(1.8), Inches(0.35))
_set(pl.paragraphs[0], "Presented by", 13, ORANGE, bold=True)
pres = textbox(s, Inches(0.9), Inches(5.95), SW - Inches(1.8), Inches(1.4))
for i, (name, sid) in enumerate([
    ("Md Wahiduzzaman Sua", "26-94088-2"),
    ("Esm-e Moula Chowdhury Abha", "26-94089-2"),
    ("Md Imtiaj Alam Sajin", "26-94090-2")]):
    p = pres.paragraphs[0] if i == 0 else pres.add_paragraph()
    p.text = f"{name}    ID: {sid}"
    r = p.runs[0]; r.font.size = Pt(15); r.font.color.rgb = WHITE; r.font.name = "Calibri"; r.font.bold = True
    p.space_after = Pt(5)

# ===== 2. DATASET =====
s = content_slide(2, "Dataset", "features + source")
bullets(s, [
    "Source: iSLEEPS, NIMHANS + IIIT-Hyderabad; published in Nature Scientific Data (2026); open via Figshare / Zenodo.",
    "100 overnight polysomnography recordings of ischemic-stroke patients; about 95,000 labelled 30-s epochs.",
    "Signal features: 6 EEG (C4:M1, C3:M2, O2:M1, O1:M2, F3/F4), 2 EOG, 3 EMG, ECG, airflow, effort, SpO2; 128 Hz down to 100 Hz.",
    "Label feature: manual AASM-2017 sleep stage per epoch (W, N1, N2, N3, REM).",
    "Clinical metadata features (63 columns): lesion side and territory, NIHSS / mRS / Barthel severity, AHI (apnea), demographics.",
], Inches(0.6), Inches(1.45), Inches(12.1), Inches(5.0), size=16.5, gap=11)
refnote(s, "Maiti S. et al., Polysomnography Dataset for Sleep Analysis in Ischemic Stroke Patients, Nature Scientific Data (2026); open access via Figshare / Zenodo.")

# ===== 3. IDENTIFY THE PROBLEM TYPE =====
s = content_slide(3, "Identify the Problem Type", "supervised")
bullets(s, [
    "A target IS given (expert AASM sleep-stage labels) -> this is SUPERVISED learning.",
    "Target / correlation we predict: the sleep stage of each 30-s epoch.",
    "The relation learned: EEG morphology (spindles, slow waves, alpha, eye-movement signatures) maps to sleep stage.",
    "Input: 4-channel EEG, 3000 samples at 100 Hz, per 30-s epoch.",
], Inches(0.6), Inches(1.45), Inches(7.5), Inches(5.6), size=17, gap=12)
# decision panel mirroring the guideline table
box(s, Inches(8.45), Inches(1.55), Inches(4.4), Inches(2.4), fill=LGRAY)
tf = textbox(s, Inches(8.7), Inches(1.7), Inches(3.9), Inches(2.2))
_set(tf.paragraphs[0], "Target not given  ->  Unsupervised / Clustering", 14, GRAY)
p = tf.add_paragraph(); _set(p, "Target IS given  ->  Supervised  (ours)", 15, NAVY, bold=True)
p = tf.add_paragraph(); _set(p, "Target = the relation we predict", 13, GRAY); p.space_before = Pt(8)

# ===== 4. WHAT TECHNIQUE ARE YOU USING? =====
s = content_slide(4, "What Technique Are You Using?", "classification")
bullets(s, [
    "Supervised CLASSIFICATION (predicts a category): 5-class sleep staging.",
    ("Not clustering (groups), not regression (forecasting), not reinforcement learning.", 1),
], Inches(0.6), Inches(1.45), Inches(12.1), Inches(1.3), size=17, gap=6)
for (lx, head, col, items) in [
    (0.55, "Deep learning", BLUE, ["CNN (per-epoch)", "CNN + BiLSTM", "DeepSleepNet", "EEG augmentation"]),
    (4.7, "Classical ML", TEAL, ["Random Forest", "Extra Trees", "XGBoost", "LightGBM"]),
    (8.85, "Hybrid", ORANGE, ["Soft-vote ensemble", "HMM (Viterbi) smoothing", "realistic stage transitions"]),
]:
    box(s, Inches(lx), Inches(2.95), Inches(3.95), Inches(3.4), fill=LGRAY)
    chip(s, Inches(lx), Inches(2.95), Inches(3.95), head, col, size=16)
    bullets(s, items, Inches(lx + 0.2), Inches(3.65), Inches(3.6), Inches(2.6), size=14, gap=9)
tf = textbox(s, Inches(0.55), Inches(6.5), Inches(12.2), Inches(0.5))
_set(tf.paragraphs[0], "All evaluated under identical subject-independent 5-fold cross-validation.", 14, GRAY, italic=True)

# ===== 5. FEATURE ENGINEERING & IMPORTANCE =====
s = content_slide(5, "Feature Engineering & Importance", "RF then XGBoost")
bullets(s, [
    "1. Feature importance: determined with RandomForest then XGBoost (gain).",
    "2. Feature engineering: per-epoch spectral band powers (delta, theta, alpha, sigma, beta), spectral entropy / edge, Hjorth, time-domain stats; 92 features, plus +/-3 epoch context (644).",
    "3. Feature selection: done last (importance-ranked).",
    ("Top features are physiologically meaningful: central delta (slow waves), occipital alpha (wake).", 1),
], Inches(0.6), Inches(1.45), Inches(6.5), Inches(5.5), size=15, gap=10)
pic(s, os.path.join(FIG, "featimp.png"), Inches(7.2), Inches(1.5), w=Inches(5.7))

# ===== 6. RESEARCH GAP =====
s = content_slide(6, "Research Gap", "+ cutting-edge models")
bullets(s, [
    "Sleep staging on healthy sleepers is saturated; the stroke-specific, lesion-aware angle is open.",
    "The base paper ignores the clinical / lesion metadata entirely (no lesion-aware error analysis).",
    "Our gap: use the staging model as an INSTRUMENT to map lesion-resolved sleep disruption, adjudicate the contested ipsilesional question, and test a severity biomarker.",
    "Cutting-edge models in scope: CNN, LSTM / BiLSTM, Transformer (used here); GNN is a natural future direction for lesion-network modeling; GAN for data augmentation.",
], Inches(0.6), Inches(1.45), Inches(12.1), Inches(5.0), size=16.5, gap=12)
refnote(s, "Spindle laterality contested: thalamic-stroke spindles, Sci. Rep. (2018); disturbed NREM laterality post-stroke (2023); thalamic-spindle study calls for lesion-symptom mapping.")

# ===== 7. BASE MODEL / BENCHMARK / BASE PAPER COMPARISON =====
s = content_slide(7, "Base Model / Benchmark / Base Paper Comparison", "journal sources")
pic(s, os.path.join(FIG, "leaderboard.png"), Inches(0.5), Inches(1.5), w=Inches(8.5))
box(s, Inches(9.4), Inches(1.6), Inches(3.5), Inches(4.9), fill=LGRAY)
bullets(s, [
    "Base paper: Maiti et al., Nature Scientific Data (2026).",
    "Published subject-independent baselines: LSTM 74.70, Transformer 67.44, CNN 61.65.",
    "Our best (full cohort, EEG+EOG+EMG ensemble+HMM): 0.742.",
    "MATCHES the published SOTA (0.742 vs 0.747, within fold noise).",
    "Journal sources: Nature, IEEE, Springer, Elsevier, MDPI, Frontiers.",
], Inches(9.6), Inches(1.75), Inches(3.15), Inches(4.6), size=12.5, gap=9)
refnote(s, "Base paper / baselines: Maiti S. et al., Nature Scientific Data (2026), subject-independent 10-fold (patient-exclusive).")

# ===== 8. RESEARCH QUESTIONS (our proposal) =====
s = content_slide(8, "Research Questions & Hypotheses", "the contribution")
qs = [
    ("RQ1", "Is sleep-microstructure disruption greater IPSILESIONALLY?  (adjudicates a contested result)", TEAL),
    ("RQ2", "Does lesion TOPOGRAPHY predict which features / stages break?  (fills a stated gap)", BLUE),
    ("RQ3", "Do signatures and model uncertainty scale with NIHSS and track mRS / Barthel?  (biomarker)", ORANGE),
    ("RQ4", "Can a subject-independent pipeline measure all this reliably?  (instrument feasibility)", GRAY),
]
y = Inches(1.55)
for tag, q, col in qs:
    box(s, Inches(0.6), y, Inches(1.1), Inches(1.2), fill=col)
    tf = textbox(s, Inches(0.6), y, Inches(1.1), Inches(1.2), MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; _set(p, tag, 19, WHITE, bold=True)
    tf2 = textbox(s, Inches(1.9), y, Inches(10.9), Inches(1.2), MSO_ANCHOR.MIDDLE)
    _set(tf2.paragraphs[0], q, 15.5, DARK)
    y = y + Inches(1.34)

# ===== 9. PRELIMINARY FINDINGS (our proposal) =====
s = content_slide(9, "Preliminary Findings", "first lesion-resolved result")
pic(s, os.path.join(FIG, "ipsilesional.png"), Inches(0.5), Inches(1.6), w=Inches(12.3))
tf = textbox(s, Inches(0.6), Inches(5.95), Inches(12.2), Inches(1.05))
_set(tf.paragraphs[0], "Robust finding: spindle-band asymmetry scales with NIHSS severity (Spearman rho = 0.41, p = 0.006, N = 43). The ipsilesional spindle reduction is directional (trend, p = 0.07).", 14.5, NAVY, bold=True)
p = tf.add_paragraph(); _set(p, "Graded by severity, not uniform, consistent with the corticothalamic literature (spindle mechanism: Fernandez & Luthi, Physiol. Rev. 2019; ipsilesional disruption: Sci. Rep. 2018, 2023). Definitive test: spindle event detection + mixed-effects + AHI control.", 12.5, GRAY)

# ===== 10. PROPOSED WORK & CONTRIBUTIONS (our proposal) =====
s = content_slide(10, "Proposed Work & Expected Contributions", "the plan")
exps = [
    ("E1", "Validate the stager on the full 100-subject cohort", "pilot done", GREEN),
    ("E2", "Ipsilesional vs contralesional microstructure (RQ1) + stats", "first result", GREEN),
    ("E3", "Lesion-topography mapping of sleep disruption (RQ2)", "proposed", GRAY),
    ("E4", "Severity / outcome biomarker from signatures + model uncertainty (RQ3)", "proposed", GRAY),
    ("E5", "Interpretability and confound control (montage, AHI, healthy reference)", "proposed", GRAY),
]
y = Inches(1.45)
for tag, txt, st, stc in exps:
    box(s, Inches(0.6), y, Inches(0.95), Inches(0.74), fill=BLUE)
    tf = textbox(s, Inches(0.6), y, Inches(0.95), Inches(0.74), MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; _set(p, tag, 15, WHITE, bold=True)
    tf2 = textbox(s, Inches(1.7), y, Inches(8.4), Inches(0.74), MSO_ANCHOR.MIDDLE)
    _set(tf2.paragraphs[0], txt, 13.5, DARK)
    box(s, Inches(10.3), y + Inches(0.12), Inches(2.5), Inches(0.5), fill=stc)
    tf3 = textbox(s, Inches(10.3), y + Inches(0.12), Inches(2.5), Inches(0.5), MSO_ANCHOR.MIDDLE)
    p = tf3.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; _set(p, st, 11, WHITE, bold=True)
    y = y + Inches(0.8)
box(s, Inches(0.6), Inches(5.6), Inches(12.2), Inches(1.35), fill=LGRAY)
tf = textbox(s, Inches(0.8), Inches(5.68), Inches(11.9), Inches(1.25), MSO_ANCHOR.MIDDLE)
_set(tf.paragraphs[0], "Contributions: first large-cohort lesion-resolved map of sleep-microstructure disruption; adjudication of the contested ipsilesional question; a candidate EEG sleep biomarker of stroke severity; reproducible pipeline.", 13, NAVY, bold=True)
p = tf.add_paragraph(); _set(p, "Target Q1 journals: Sleep, Journal of Sleep Research, NeuroImage Clinical, IEEE J-BHI, Computers in Biology and Medicine.", 12.5, GRAY)

prs.save(OUT)
print("saved ->", OUT, "| slides:", len(prs.slides._sldIdLst))
