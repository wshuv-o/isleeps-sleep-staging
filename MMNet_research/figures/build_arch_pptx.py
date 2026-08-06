"""Assemble the MM-Net architecture figure as a native, editable PowerPoint slide, using the
real data mini-plots (arch_assets/) and the 3D encoder blocks. A matplotlib preview is rendered
from the SAME coordinates so the layout can be verified here (PowerPoint cannot be rendered).
Outputs: MM_Net_architecture.pptx  +  arch_pptx_preview.png"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
import matplotlib; matplotlib.use("Agg")
import matplotlib.pyplot as plt, matplotlib.image as mpimg
import matplotlib.patches as mp
from matplotlib.patches import FancyArrowPatch

HERE = os.path.dirname(os.path.abspath(__file__)); A = os.path.join(HERE, "arch_assets")
def rgb(h):
    h=h.lstrip('#')
    if len(h)==3: h=''.join(c*2 for c in h)
    return RGBColor(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))
GREEN,ORANGE,BLUE,YELLOW,PURPLE,RED,GREY = "#e8f5ec","#fdecdd","#e7effb","#fdf6da","#efe7fb","#fde7e7","#f4f5f7"
GST,OST,BST,YST,PST,RST = "#3f8f5a","#c9803a","#3a6ea5","#b6902a","#6b46c1","#c53030"

EL = []  # shared layout model for preview
def _box(x,y,w,h,text,fill,line,fs,bold,rounded,dashed):
    EL.append(("box",x,y,w,h,text,fill,line,fs,bold,rounded,dashed))
def _img(x,y,w,h,name): EL.append(("img",x,y,w,h,name))
def _txt(x,y,w,h,text,fs,bold,color): EL.append(("txt",x,y,w,h,text,fs,bold,color))
def _circ(x,y,d,fill,line): EL.append(("circ",x,y,d,fill,line))
def _arr(x1,y1,x2,y2,color,dashed): EL.append(("arr",x1,y1,x2,y2,color,dashed))

# ---------------- layout (inches) on a 13.33 x 7.5 slide ----------------
def box(x,y,w,h,text="",fill="#ffffff",line="#333333",fs=11,bold=True,rounded=True,dashed=False):
    _box(x,y,w,h,text,fill,line,fs,bold,rounded,dashed); return (x,y,w,h)
def img(x,y,w,h,name): _img(x,y,w,h,name)
def txt(x,y,w,h,text,fs=10,bold=False,color="#222222"): _txt(x,y,w,h,text,fs,bold,color)
def neuroncol(cx,cy,n,color,line,label):   # vertical pill of small circles
    d=0.11; box(cx-0.16,cy-0.05,0.32,n*d+0.1,"",color,line,8,False,True,False)
    for i in range(n): _circ(cx-d/2, cy+i*d, d, color, line)
    txt(cx-0.4,cy+n*d+0.06,0.8,0.18,label,8,True,line)
def lstmchain(x,y,w):
    cell=0.5; gap=(w-5*cell)/4
    for r,(arrow_c,yoff,tag) in enumerate([("#3f8f5a",0,"fwd"),("#c9803a",0.62,"bwd")]):
        for i in range(5):
            bx=x+i*(cell+gap)
            box(bx,y+yoff,cell,0.34,"LSTM","#efe7fb",PST,7,False,True,False)
            if i<4:
                if r==0: _arr(bx+cell,y+yoff+0.17,bx+cell+gap,y+yoff+0.17,arrow_c,False)
                else: _arr(bx+cell+gap,y+yoff+0.17,bx+cell,y+yoff+0.17,arrow_c,False)
    for i in range(5): txt(x+i*(cell+gap),y+1.02,cell,0.16,f"x{i+1}",8,False,"#555")

# EEG stream
b_eegrec=box(0.15,0.30,1.85,1.45,"EEG Recordings",GREEN,GST,10); img(0.28,0.62,1.6,1.05,"eeg_traces.png")
b_eegfe =box(2.15,0.30,1.95,1.45,"EEG Feature Extraction\n(band power, spindle, Hjorth)",GREEN,GST,9); img(2.55,0.78,1.15,0.9,"bandpower.png")
b_eegen =box(4.25,0.45,1.35,1.15,"EEG Encoder\n(phi_eeg, FeatMLP)",GREEN,GST,9); img(4.30,0.78,1.25,0.72,"eeg_mlp_3d.png")
neuroncol(5.95,0.55,6,GREEN,GST,"z_e 128")

# Cardio stream
b_carsig=box(0.15,2.05,1.85,1.45,"Cardio Signals",ORANGE,OST,10); img(0.28,2.37,1.6,1.05,"cardio_signals.png")
b_carfe =box(2.15,2.05,1.95,1.45,"Cardio Feature Extraction\n(SpO2, effort, HRV, airflow)",ORANGE,OST,9); img(2.65,2.55,1.0,0.9,"hrv.png")
b_caren =box(4.25,2.20,1.35,1.15,"Cardio Encoder\n(phi_car, FeatMLP)",ORANGE,OST,9); img(4.30,2.53,1.25,0.72,"cardio_mlp_3d.png")
neuroncol(5.95,2.35,5,ORANGE,OST,"z_c 64")

# Cross-modal fusion
b_fus=box(6.45,0.95,2.35,1.95,"Cross-modal Fusion\n(2 tokens + attention)\nz in R^128",BLUE,BST,10)
for k in range(4): img(6.60+k*0.52,1.55,0.46,0.46,f"attn_head{k+1}.png")
txt(6.60,2.02,2.05,0.16,"head 1   head 2   head 3   head 4",7,False,BST)
neuroncol(8.55,1.30,6,BLUE,BST,"z 128")

# BiLSTM
b_lstm=box(2.7,3.75,4.6,1.5,"BiLSTM  (2 layers, bidirectional)   context L = 20   h_t in R^256",YELLOW,YST,10)
lstmchain(2.95,4.08,4.1)

# heads
b_stage=box(1.55,5.45,3.1,1.75,"Staging Head\ny_stg = softmax(W_s h_t) + HMM Viterbi",PURPLE,PST,9); img(1.62,5.95,1.9,1.15,"hypnogram_mini.png")
for i,l in enumerate(["Wake","REM","N1","N2","N3"]): _circ(3.75,5.98+i*0.22,0.16,PURPLE,PST); txt(3.95,5.96+i*0.22,0.6,0.18,l,7,False,"#333")
b_resp=box(4.95,5.45,3.1,1.75,"Respiratory Head (Apnea)\ny_apn = sigmoid(W_r [h_t ; c_t])",RED,RST,9); img(5.02,6.05,1.9,0.95,"apnea_wave.png")
for i,l in enumerate(["Apnea","Non-Apnea"]): _circ(7.15,6.2+i*0.30,0.18,RED,RST); txt(7.36,6.18+i*0.30,0.85,0.2,l,7,False,"#333")

# joint objective
box(2.4,7.15,5.4,0.3,"Joint objective:  L = CE_stg(y_stg,target) + lambda BCE_apn(y_apn,target),  lambda=1","#ffffff","#333333",9,True,False,True)

# detail panel 1 (FeatMLP)
box(9.05,0.30,2.05,3.0,"Feature Encoder (FeatMLP)\ne = GELU(LN(W2 GELU(LN(W1 x))))",GREEN,GST,9)
for j,t in enumerate(["Linear W1","LayerNorm","GELU","Dropout (0.3)","Linear W2","LayerNorm","GELU"]):
    box(9.30,0.86+j*0.31,1.55,0.26,t,"#d7ecdd",GST,8,False,True,False)
txt(9.05,3.05,2.05,0.18,"EEG: 188 to 128   -   Cardio: 14 to 64",8,False,GST)

# detail panel 2 (Attention)
box(9.05,3.45,2.05,3.05,"Cross-modal Fusion (Attention)",BLUE,BST,9)
for j,t in enumerate(["Tokenize T=[W_e z_e;W_c z_c]+M","Compute Q, K, V","4 heads - concat - W_o","Residual + LayerNorm","Feed-Forward W_f","fuse -> z in R^128"]):
    box(9.22,3.80+j*0.42,1.72,0.32,t,"#dbe6f7",BST,8,False,True,False)

# legend
box(0.15,7.15,2.1,0.3,"data flow ->   direct c_t (dashed)   attention   LSTM chain","#ffffff","#999999",8,False,False,False)

# arrows connect box edge-to-edge (cleaner routing)
def edgepts(b1,b2):
    x1,y1,w1,h1=b1; x2,y2,w2,h2=b2; c1=(x1+w1/2,y1+h1/2); c2=(x2+w2/2,y2+h2/2)
    dx=c2[0]-c1[0]; dy=c2[1]-c1[1]
    if abs(dx)>=abs(dy):
        return ((x1+w1,c1[1]),(x2,c2[1])) if dx>0 else ((x1,c1[1]),(x2+w2,c2[1]))
    return ((c1[0],y1+h1),(c2[0],y2)) if dy>0 else ((c1[0],y1),(c2[0],y2+h2))
def arrow(b1,b2,color="#333",dashed=False):
    (px1,py1),(px2,py2)=edgepts(b1,b2); _arr(px1,py1,px2,py2,color,dashed)
arrow(b_eegrec,b_eegfe); arrow(b_eegfe,b_eegen); arrow(b_carsig,b_carfe); arrow(b_carfe,b_caren)
arrow(b_eegen,b_fus); arrow(b_caren,b_fus); arrow(b_fus,b_lstm); arrow(b_lstm,b_stage); arrow(b_lstm,b_resp)
arrow(b_caren,b_resp,"#c9803a",True)

# ================= emit PowerPoint =================
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
sl=prs.slides.add_slide(prs.slide_layouts[6])
def setfill(sh,fill,line):
    sh.fill.solid(); sh.fill.fore_color.rgb=rgb(fill); sh.line.color.rgb=rgb(line); sh.line.width=Pt(1.2); sh.shadow.inherit=False
def settext(sh,text,fs,bold,color,anchor_top=True):
    tf=sh.text_frame; tf.word_wrap=True; tf.vertical_anchor=MSO_ANCHOR.TOP if anchor_top else MSO_ANCHOR.MIDDLE
    tf.margin_left=Pt(3); tf.margin_right=Pt(3); tf.margin_top=Pt(2); tf.margin_bottom=Pt(2)
    for k,line in enumerate(text.split("\n")):
        p=tf.paragraphs[0] if k==0 else tf.add_paragraph(); p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=line; r.font.size=Pt(fs); r.font.bold=bold; r.font.color.rgb=rgb(color)
def add_arrow(x1,y1,x2,y2,color,dashed):
    cn=sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    cn.line.color.rgb=rgb(color); cn.line.width=Pt(1.4); ln=cn.line._get_or_add_ln()
    te=ln.makeelement(qn('a:tailEnd'),{'type':'triangle','w':'med','len':'med'}); ln.append(te)
    if dashed:
        d=ln.makeelement(qn('a:prstDash'),{'val':'dash'}); ln.append(d)
for e in EL:
    if e[0]=="box":
        _,x,y,w,h,text,fill,line,fs,bold,rounded,dashed=e
        shp=MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
        sh=sl.shapes.add_shape(shp,Inches(x),Inches(y),Inches(w),Inches(h)); setfill(sh,fill,line)
        if dashed: sh.line._get_or_add_ln().append(sh.line._get_or_add_ln().makeelement(qn('a:prstDash'),{'val':'dash'}))
        if text: settext(sh,text,fs,bold,line if bold else "#333333")
    elif e[0]=="img":
        _,x,y,w,h,name=e; p=os.path.join(A,name)
        if os.path.exists(p): sl.shapes.add_picture(p,Inches(x),Inches(y),Inches(w),Inches(h))
    elif e[0]=="txt":
        _,x,y,w,h,text,fs,bold,color=e
        tb=sl.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); settext(tb,text,fs,bold,color)
    elif e[0]=="circ":
        _,x,y,d,fill,line=e; sh=sl.shapes.add_shape(MSO_SHAPE.OVAL,Inches(x),Inches(y),Inches(d),Inches(d)); setfill(sh,fill,line)
    elif e[0]=="arr":
        _,x1,y1,x2,y2,color,dashed=e; add_arrow(x1,y1,x2,y2,color,dashed)
prs.save(os.path.join(HERE,"MM_Net_architecture.pptx"))

# ================= emit matplotlib preview (same coords) =================
fig,ax=plt.subplots(figsize=(13.33,7.5)); ax.set_xlim(0,13.333); ax.set_ylim(0,7.5); ax.invert_yaxis(); ax.axis("off")
for e in EL:
    if e[0]=="box":
        _,x,y,w,h,text,fill,line,fs,bold,rounded,dashed=e
        ax.add_patch(mp.FancyBboxPatch((x,y),w,h,boxstyle="round,pad=0,rounding_size=0.05" if rounded else "square,pad=0",
            linewidth=1.1,edgecolor=line,facecolor=fill,linestyle="--" if dashed else "-"))
        if text: ax.text(x+w/2,y+0.04,text,ha="center",va="top",fontsize=fs*0.8,fontweight="bold" if bold else "normal",color=line if bold else "#333")
    elif e[0]=="txt":
        _,x,y,w,h,text,fs,bold,color=e; ax.text(x+w/2,y+h/2,text,ha="center",va="center",fontsize=fs*0.8,fontweight="bold" if bold else "normal",color=color)
    elif e[0]=="img":
        _,x,y,w,h,name=e; p=os.path.join(A,name)
        if os.path.exists(p): ax.imshow(mpimg.imread(p),extent=(x,x+w,y+h,y),aspect="auto",zorder=5)
    elif e[0]=="circ":
        _,x,y,d,fill,line=e; ax.add_patch(mp.Circle((x+d/2,y+d/2),d/2,facecolor=fill,edgecolor=line,lw=0.8,zorder=6))
    elif e[0]=="arr":
        _,x1,y1,x2,y2,color,dashed=e; ax.add_patch(FancyArrowPatch((x1,y1),(x2,y2),arrowstyle="-|>",mutation_scale=10,color=color,lw=1.1,linestyle="--" if dashed else "-",zorder=2))
fig.savefig(os.path.join(HERE,"arch_pptx_preview.png"),dpi=100,bbox_inches="tight")
print("wrote MM_Net_architecture.pptx + arch_pptx_preview.png")
